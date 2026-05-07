#!/usr/bin/env python3
"""
PPT 设计思维自动学习引擎
========================
六层分析：文字提取 → 图片提取 → 图文关系 → 叙事结构 → 概念方案 → 思维模式

核心方法论：点 → 线 → 面 → 空间

用法：
    python ppt_learner.py              # 扫描 ppt/ 下所有新文件
    python ppt_learner.py --all        # 重新处理所有文件
    python ppt_learner.py "path.pptx"  # 处理指定文件
    python ppt_learner.py --watch      # 持续监听模式
"""

import json, os, sys, re, hashlib, shutil, time, zipfile
from pathlib import Path
from datetime import datetime
from xml.etree import ElementTree as ET

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
except ImportError:
    print("请先安装 python-pptx: pip install python-pptx")
    sys.exit(1)

# ═══════════════════════════════════════════════
# 配置
# ═══════════════════════════════════════════════

ROOT = Path("f:/code/ARK Design")
PPT_DIR = ROOT / "ppt"
KNOWLEDGE_DIR = ROOT / ".ppt-knowledge"
PROCESSED_LOG = KNOWLEDGE_DIR / "processed.json"
INDEX_FILE = KNOWLEDGE_DIR / "index.md"

MASTERS_DIR = KNOWLEDGE_DIR / "masters"
PROMPT_PATTERNS_DIR = KNOWLEDGE_DIR / "prompt-patterns"
NARRATIVE_ARCS_DIR = KNOWLEDGE_DIR / "narrative-arcs"
OPENINGS_DIR = KNOWLEDGE_DIR / "openings"
SPEECH_RHYTHMS_DIR = KNOWLEDGE_DIR / "speech-rhythms"
CLOSINGS_DIR = KNOWLEDGE_DIR / "closings"
CONCEPT_PATTERNS_DIR = KNOWLEDGE_DIR / "concept-patterns"


def sanitize_name(name):
    """清理文件名：去掉尾部空格和点（Windows兼容），限制长度"""
    name = name.rstrip('. ')
    if not name:
        name = "unnamed"
    return name


def ensure_dirs():
    for d in [KNOWLEDGE_DIR, MASTERS_DIR, PROMPT_PATTERNS_DIR,
              NARRATIVE_ARCS_DIR, OPENINGS_DIR, SPEECH_RHYTHMS_DIR,
              CLOSINGS_DIR, CONCEPT_PATTERNS_DIR]:
        d.mkdir(parents=True, exist_ok=True)


# ═══════════════════════════════════════════════
# 已处理文件跟踪
# ═══════════════════════════════════════════════

def load_processed():
    if PROCESSED_LOG.exists():
        return json.loads(PROCESSED_LOG.read_text(encoding='utf-8'))
    return {}

def save_processed(processed):
    PROCESSED_LOG.write_text(json.dumps(processed, ensure_ascii=False, indent=2), encoding='utf-8')

def file_fingerprint(path):
    """计算文件特征：大小 + 修改时间，用于增量检测"""
    st = os.stat(path)
    return f"{st.st_size}_{int(st.st_mtime)}"

def is_new_file(path, processed):
    name = path.name
    fp = file_fingerprint(path)
    return name not in processed or processed[name] != fp


# ═══════════════════════════════════════════════
# PPT 解析核心
# ═══════════════════════════════════════════════

def extract_pptx(pptx_path):
    """
    提取PPTX的所有内容，返回结构化数据：
    {
        "file": { "name": ..., "pages": ..., "size": ... },
        "slides": [
            {
                "index": 1,
                "texts": [ { "type": "title|body|note", "text": "...", "level": 0 } ],
                "images": [ { "filename": "...", "width": ..., "height": ... } ],
                "notes": "..."
            }
        ]
    }
    """
    pptx_path = Path(pptx_path)  # 兼容 str 和 Path
    prs = Presentation(str(pptx_path))
    result = {
        "file": {
            "name": pptx_path.name,
            "size": os.path.getsize(pptx_path),
            "pages": len(prs.slides),
        },
        "slides": []
    }

    for idx, slide in enumerate(prs.slides, 1):
        slide_data = {"index": idx, "texts": [], "images": [], "notes": ""}

        # 提取所有文字
        for shape in slide.shapes:
            # 正文文本框
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        # 判断近似标题：字体较大或位置靠前
                        is_heading = False
                        if para.runs:
                            font_size = para.runs[0].font.size
                            if font_size and font_size >= Pt(18):
                                is_heading = True
                        slide_data["texts"].append({
                            "type": "heading" if is_heading else "body",
                            "text": text,
                            "level": para.level if hasattr(para, 'level') else 0
                        })

            # 图片信息（占位）
            if shape.shape_type == 13:  # picture
                try:
                    img = shape.image
                    slide_data["images"].append({
                        "filename": f"slide-{idx:02d}-img{len(slide_data['images'])+1}.{img.content_type.split('/')[-1]}",
                        "width": shape.width,
                        "height": shape.height,
                        "left": shape.left,
                        "top": shape.top,
                        "content_type": img.content_type
                    })
                except:
                    pass

        # 提取备注
            if shape.has_text_frame:
                pass  # 已在上方处理
        # 备注在 slide.notes_slide
        try:
            notes_slide = slide.notes_slide
            notes_text = notes_slide.notes_text_frame.text.strip()
            slide_data["notes"] = notes_text
        except:
            pass

        result["slides"].append(slide_data)

    return result


def extract_images_from_pptx(pptx_path, output_dir):
    """
    直接从PPTX（ZIP）中解压图片到 output_dir
    返回已提取的图片文件列表
    """
    extracted = []
    pptx_path = str(pptx_path)

    # 图片在PPTX压缩包中的路径: ppt/media/imageN.png
    image_extensions = {'.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff', '.webp', '.svg'}

    try:
        with zipfile.ZipFile(pptx_path, 'r') as z:
            for name in z.namelist():
                ext = Path(name).suffix.lower()
                if ext in image_extensions and name.startswith('ppt/media/'):
                    # 保持原始文件名
                    base = Path(name).name
                    out_path = output_dir / base
                    # 避免重名
                    if out_path.exists():
                        base = f"img_{len(extracted)}_{base}"
                        out_path = output_dir / base
                    with z.open(name) as src:
                        out_path.write_bytes(src.read())
                    extracted.append(base)
    except Exception as e:
        print(f"  [警告] 图片解压失败: {e}")

    return extracted


# ═══════════════════════════════════════════════
# 分析：图文关系
# ═══════════════════════════════════════════════

def analyze_image_text_relationship(slide_data):
    """
    分析一页内图文关系：识别图文搭配逻辑
    注：真正的语义分析需要 Claude 处理，这里做结构化基础分析
    """
    texts = [t["text"] for t in slide_data["texts"]]
    img_count = len(slide_data["images"])
    notes = slide_data["notes"]

    # 基础分析：图文密度、位置关系描述
    total_chars = sum(len(t) for t in texts)
    relationship = {
        "image_count": img_count,
        "text_char_count": total_chars,
        "has_notes": bool(notes),
        "text_density": "high" if total_chars > 200 else ("medium" if total_chars > 80 else "low"),
    }
    return relationship


# ═══════════════════════════════════════════════
# 分析：概念方案（点→线→面→空间）
# ═══════════════════════════════════════════════

def analyze_concept_scheme(slide_texts, all_slides):
    """
    分析概念方案的点→线→面→空间结构。
    注意：这个分析是初步的结构化分析，真正的深度分析
    需要 Claude 的语义理解能力来完成。
    """
    # 收集所有文本
    all_text = "\n".join(s["text"] for s in slide_texts)

    # 识别概念原点关键词
    origin_keywords = {
        "文化": "文化符号",
        "传统": "文化符号",
        "当地": "文化符号",
        "自然": "自然元素",
        "光影": "自然元素",
        "光线": "自然元素",
        "材质": "材质肌理",
        "肌理": "材质肌理",
        "问题": "问题痛点",
        "痛点": "问题痛点",
        "生活方式": "生活方式",
        "体验": "体验场景",
    }

    detected_origins = set()
    for keyword, category in origin_keywords.items():
        if keyword in all_text:
            detected_origins.add(category)

    # 识别设计语言词汇
    lang_keywords = {
        "材质": False, "材料": False,
        "色彩": False, "颜色": False,
        "灯光": False, "照明": False,
        "家具": False, "软装": False,
        "空间": False, "布局": False,
    }
    for kw in lang_keywords:
        if kw in all_text:
            lang_keywords[kw] = True

    return {
        "detected_origins": list(detected_origins),
        "design_language_hits": lang_keywords,
        "concept_slide_count": len(all_slides) if "概念" in all_text or "方案" in all_text else None,
    }


# ═══════════════════════════════════════════════
# 分析：叙事结构
# ═══════════════════════════════════════════════

def analyze_narrative(slides):
    """
    初步分析叙事结构：识别故事弧线类型、节奏
    """
    total = len(slides)
    if total == 0:
        return {"structure_type": "unknown", "slide_count": 0}

    # 收集每页文本量作为节奏指标
    page_densities = []
    for s in slides:
        char_count = sum(len(t["text"]) for t in s["texts"])
        page_densities.append(char_count)

    # 有备注的页面数量
    notes_count = sum(1 for s in slides if s["notes"])

    # 推测结构类型（基于页面数量和分布）
    if total <= 8:
        structure_type = "简洁型"
    elif total <= 15:
        structure_type = "标准提案型"
    else:
        structure_type = "详细报告型"

    return {
        "structure_type": structure_type,
        "slide_count": total,
        "notes_ratio": f"{notes_count}/{total}",
        "density_curve": "→".join(str(d) for d in page_densities[:12]) + ("..." if total > 12 else ""),
    }


# ═══════════════════════════════════════════════
# 输出：生成结构化知识
# ═══════════════════════════════════════════════

def generate_extract_md(pptx_data, slide_images):
    """生成原文提取 markdown"""
    lines = []
    lines.append(f"# {pptx_data['file']['name']}")
    lines.append(f"")
    lines.append(f"- 页数：{pptx_data['file']['pages']}")
    lines.append(f"- 大小：{pptx_data['file']['size'] // 1024}KB")
    lines.append(f"- 处理时间：{datetime.now().strftime('%Y-%m-%d %H:%M')}")
    lines.append(f"")
    lines.append("---")
    lines.append("")

    for slide in pptx_data["slides"]:
        lines.append(f"## 第 {slide['index']} 页")
        lines.append("")
        for t in slide["texts"]:
            prefix = "###" if t["type"] == "heading" else ""
            if prefix:
                lines.append(f"{prefix} {t['text']}")
            else:
                lines.append(t["text"])
            lines.append("")

        # 本页图片
        page_imgs = [img for img in slide_images if f"slide-{slide['index']:02d}" in img]
        if page_imgs:
            lines.append(f"*本页关联图片：{'、'.join(page_imgs)}*")
            lines.append("")

        # 备注
        if slide["notes"]:
            lines.append("> **演讲备注**")
            for note_line in slide["notes"].split("\n"):
                if note_line.strip():
                    lines.append(f"> {note_line.strip()}")
            lines.append("")

        lines.append("---")
        lines.append("")

    return "\n".join(lines)


def generate_analysis_md(pptx_data):
    """生成图文关系分析 markdown"""
    lines = []
    lines.append(f"# 图文关系分析：{pptx_data['file']['name']}")
    lines.append("")
    lines.append(f"总页数：{pptx_data['file']['pages']}")
    lines.append("")

    total_images = sum(len(s["images"]) for s in pptx_data["slides"])
    total_notes = sum(1 for s in pptx_data["slides"] if s["notes"])

    lines.append("## 统计概览")
    lines.append(f"- 总图片数：{total_images}")
    lines.append(f"- 有备注的页数：{total_notes} / {pptx_data['file']['pages']}")
    lines.append("")

    lines.append("## 逐页图文概览")
    lines.append("")
    lines.append("| 页面 | 文字量 | 图片数 | 有备注 | 密度 |")
    lines.append("|------|--------|--------|--------|------|")
    for s in pptx_data["slides"]:
        rel = analyze_image_text_relationship(s)
        lines.append(f"| {s['index']} | {rel['text_char_count']}字 | {rel['image_count']} | {'Y' if rel['has_notes'] else ''} | {rel['text_density']} |")
    lines.append("")

    return "\n".join(lines)


def generate_concept_md(pptx_data, pptx_path):
    """生成概念方案分析 markdown（点→线→面→空间框架）"""
    all_texts = [t for s in pptx_data["slides"] for t in s["texts"]]
    concept = analyze_concept_scheme(all_texts, pptx_data["slides"])

    lines = []
    lines.append(f"# 概念方案分析：{pptx_data['file']['name']}")
    lines.append("")
    lines.append(f"## 点：概念原点")
    lines.append("")
    if concept["detected_origins"]:
        for o in concept["detected_origins"]:
            lines.append(f"- {o}")
    else:
        lines.append("*待Claude深度识别*")
    lines.append("")

    lines.append(f"## 线：推演逻辑")
    lines.append("")
    lines.append("*待Claude深度分析*")
    lines.append("")

    lines.append(f"## 面：设计语言体系")
    lines.append("")
    hit_keys = [k for k, v in concept["design_language_hits"].items() if v]
    if hit_keys:
        lines.append("检测到的设计语言维度：")
        for k in hit_keys:
            lines.append(f"- {k}")
    else:
        lines.append("*待Claude深度识别*")
    lines.append("")

    lines.append(f"## 空间：完整体验")
    lines.append("")
    lines.append("*待Claude深度分析*")
    lines.append("")

    return "\n".join(lines)


def generate_narrative_md(pptx_data):
    """生成叙事结构分析 markdown"""
    narrative = analyze_narrative(pptx_data["slides"])

    lines = []
    lines.append(f"# 叙事结构分析：{pptx_data['file']['name']}")
    lines.append("")
    lines.append(f"结构类型：{narrative['structure_type']}")
    lines.append(f"页数：{narrative['slide_count']}")
    lines.append(f"有备注的页面：{narrative['notes_ratio']}")
    lines.append("")
    lines.append("## 节奏曲线（每页文字量）")
    lines.append("")
    lines.append(f"```")
    lines.append(narrative["density_curve"])
    lines.append(f"```")
    lines.append("")
    lines.append("## 演讲备注摘要")
    lines.append("")
    notes_found = False
    for s in pptx_data["slides"]:
        if s["notes"]:
            notes_found = True
            first_line = s["notes"].strip().split("\n")[0][:120]
            lines.append(f"- 第{s['index']}页备注：{first_line}...")
    if not notes_found:
        lines.append("*本PPT无演讲备注*")
    lines.append("")
    lines.append("## 叙事分析")
    lines.append("")
    lines.append("*待Claude深度分析叙事弧线、开场技法、节奏控制、收尾艺术*")
    lines.append("")

    return "\n".join(lines)


def generate_patterns_md(pptx_data):
    """生成可复用模式汇总 markdown"""
    all_text = []
    for s in pptx_data["slides"]:
        for t in s["texts"]:
            all_text.append(t["text"])

    # 提取金句风格（简短有力的句子）
    gold_sentences = []
    for text in all_text:
        for s in re.split(r'[。！？\n]', text):
            s = s.strip()
            # 15-60字，可能包含关键修辞
            if 10 < len(s) < 80 and any(kw in s for kw in ["不是", "而是", "最", "本质", "设计", "空间", "体验", "在于"]):
                gold_sentences.append(s)

    lines = []
    lines.append(f"# 可复用模式：{pptx_data['file']['name']}")
    lines.append("")
    lines.append("## 可能的设计金句")
    lines.append("")
    if gold_sentences:
        for s in gold_sentences[:10]:
            lines.append(f"- 「{s}」")
    else:
        lines.append("*待Claude深度提取*")
    lines.append("")

    lines.append("## 图文搭配模式")
    lines.append("")
    lines.append("*待Claude深度分析*")
    lines.append("")

    lines.append("## 概念推演模式")
    lines.append("")
    lines.append("*待Claude深度分析*")
    lines.append("")

    return "\n".join(lines)


# ═══════════════════════════════════════════════
# 主处理流程
# ═══════════════════════════════════════════════

def process_pptx(pptx_path):
    """处理单个PPT文件，输出到知识库"""
    safe_name = pptx_path.name.encode('utf-8', errors='replace').decode('utf-8', errors='replace')
    print(f"\n{'='*60}")
    print(f"学习: {safe_name}")
    print(f"{'='*60}")

    # 1. 解析PPTX
    pptx_data = extract_pptx(str(pptx_path))
    print(f"  页数: {pptx_data['file']['pages']}")
    print(f"  开始分析...")

    # 2. 创建大师目录（处理长文件名和特殊字符）
    master_name = sanitize_name(pptx_path.stem)
    master_dir = MASTERS_DIR / master_name
    slides_dir = master_dir / "slides"
    try:
        slides_dir.mkdir(parents=True, exist_ok=True)
    except OSError as e:
        print(f"  [警告] 无法创建目录，使用短名称: {e}")
        # 尝试使用短名称（截断+hash）
        short_name = master_name[:60] + "_" + hashlib.md5(master_name.encode()).hexdigest()[:8]
        master_dir = MASTERS_DIR / short_name
        slides_dir = master_dir / "slides"
        slides_dir.mkdir(parents=True, exist_ok=True)
        print(f"  使用短名称: {short_name}")

    # 3. 提取图片
    print(f"  提取图片...")
    extracted_images = extract_images_from_pptx(pptx_path, slides_dir)
    print(f"  提取了 {len(extracted_images)} 张图片")

    # 4. 按页建立图文映射
    for slide in pptx_data["slides"]:
        # 找到本页最可能的图片
        page_images = [img for img in extracted_images if f"slide-{slide['index']:02d}" in img]
        if not page_images:
            # 如果没找到按页命名的，就尝试位置匹配
            pass

    # 5. 生成分析文档
    print(f"  生成分析文档...")

    (master_dir / "extract.md").write_text(
        generate_extract_md(pptx_data, extracted_images), encoding='utf-8')
    print(f"    [OK] extract.md")

    (master_dir / "analysis.md").write_text(
        generate_analysis_md(pptx_data), encoding='utf-8')
    print(f"    [OK] analysis.md")

    (master_dir / "narrative.md").write_text(
        generate_narrative_md(pptx_data), encoding='utf-8')
    print(f"    [OK] narrative.md")

    (master_dir / "concept.md").write_text(
        generate_concept_md(pptx_data, pptx_path), encoding='utf-8')
    print(f"    [OK] concept.md")

    (master_dir / "patterns.md").write_text(
        generate_patterns_md(pptx_data), encoding='utf-8')
    print(f"    [OK] patterns.md")

    # 6. 保存原始数据的JSON（供Claude深度分析参考）
    json_path = master_dir / "raw.json"
    # 简化JSON，只保留文本和图片信息
    raw_simple = {
        "source": pptx_data["file"]["name"],
        "page_count": pptx_data["file"]["pages"],
        "slides": [{
            "index": s["index"],
            "texts": [{"type": t["type"], "text": t["text"]} for t in s["texts"]],
            "images": [i["filename"] for i in s["images"]],
            "notes": s["notes"][:500]  # 截断以防过大
        } for s in pptx_data["slides"]]
    }
    json_path.write_text(json.dumps(raw_simple, ensure_ascii=False, indent=2), encoding='utf-8')
    print(f"    [OK] raw.json")

    print(f"  学习完成: {master_name}")
    return True


def scan_and_learn():
    """扫描 ppt/ 目录，处理所有新文件"""
    ensure_dirs()
    processed = load_processed()
    new_files = []

    if not PPT_DIR.exists():
        print(f"PPT目录不存在: {PPT_DIR}")
        return

    # 收集所有 pptx 文件
    for f in sorted(PPT_DIR.glob("*.pptx")):
        if f.is_file() and is_new_file(f, processed):
            new_files.append(f)

    if not new_files:
        print(f"没有新PPT文件需要学习。（共 {len(processed)} 个已学习）")
        return

    print(f"发现 {len(new_files)} 个新PPT文件，开始学习...")

    success_count = 0
    for f in new_files:
        try:
            process_pptx(f)
            processed[f.name] = file_fingerprint(f)
            save_processed(processed)
            success_count += 1
        except Exception as e:
            print(f"  [错误] 处理失败: {e}")
            import traceback
            traceback.print_exc()

    # 更新索引
    update_index()

    print(f"\n{'='*60}")
    print(f"学习完成：成功 {success_count} / {len(new_files)}")
    print(f"知识库位置：{KNOWLEDGE_DIR}")


def process_single(path):
    """处理指定文件"""
    ensure_dirs()
    path = Path(path)
    if not path.exists():
        print(f"文件不存在: {path}")
        return

    if path.suffix.lower() not in ('.pptx',):
        print(f"不支持的文件格式: {path.suffix}（仅支持 .pptx）")
        return

    processed = load_processed()
    try:
        process_pptx(path)
        processed[path.name] = file_fingerprint(path)
        save_processed(processed)
        update_index()
    except Exception as e:
        print(f"处理失败: {e}")
        import traceback
        traceback.print_exc()


def update_index():
    """更新索引文件"""
    processed = load_processed()

    lines = []
    lines.append("# PPT 设计思维知识库")
    lines.append("")
    lines.append(f"最后更新：{datetime.now().strftime('%Y-%m-%d %H:%M')}")
    lines.append("")
    lines.append(f"已学习PPT：{len(processed)} 份")
    lines.append("")
    lines.append("## 已学习清单")
    lines.append("")
    lines.append("| 文件名 | 学习时间 |")
    lines.append("|--------|----------|")

    for name in sorted(processed.keys()):
        fp = processed[name]
        # 从 fingerprint 无法得知时间，用索引更新时间
        lines.append(f"| {name} | {datetime.now().strftime('%Y-%m-%d')} |")

    lines.append("")
    lines.append("---")
    lines.append("")
    lines.append("## 知识库目录")
    lines.append("")
    lines.append("- [masters/](masters/) — 大师学习笔记")
    lines.append("- [concept-patterns/](concept-patterns/) — ★ 概念方案模式库")
    lines.append("- [narrative-arcs/](narrative-arcs/) — 叙事框架库")
    lines.append("- [prompt-patterns/](prompt-patterns/) — 图文/生图模式库")
    lines.append("- [openings/](openings/) — 开场技法库")
    lines.append("- [speech-rhythms/](speech-rhythms/) — 演讲节奏模板")
    lines.append("- [closings/](closings/) — 收尾技法库")

    INDEX_FILE.write_text("\n".join(lines), encoding='utf-8')
    print(f"\n索引已更新: {INDEX_FILE}")


# ═══════════════════════════════════════════════
# 命令行入口
# ═══════════════════════════════════════════════

if __name__ == "__main__":
    import sys

    if len(sys.argv) > 1:
        arg = sys.argv[1]
        if arg == "--all":
            # 强制全部重新处理
            if PROCESSED_LOG.exists():
                PROCESSED_LOG.unlink()
            scan_and_learn()
        elif arg == "--watch":
            print("监听模式：按 Ctrl+C 停止")
            try:
                while True:
                    scan_and_learn()
                    print("\n等待 60 秒后再次检查...")
                    time.sleep(60)
            except KeyboardInterrupt:
                print("\n监听已停止")
        elif arg == "--help" or arg == "-h":
            print(__doc__)
        else:
            # 处理指定文件
            process_single(arg)
    else:
        # 默认：扫描新文件
        scan_and_learn()
