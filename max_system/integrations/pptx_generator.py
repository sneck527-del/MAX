"""将 ARK Design narrator 输出的 pages 数组转换为 PPTX 文件。

轻量级 generator，不依赖 ARK — 仅读取 result.json 中的 pages 数组即可。
"""

from __future__ import annotations

import logging
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

logger = logging.getLogger(__name__)

# ── 7 个 Section 名称 ──────────────────────────────────
SECTION_NAMES = {
    1: "项目概况",
    2: "概念推演",
    3: "空间叙事",
    4: "效果图展示",
    5: "灯光与材质",
    6: "风险与预算",
    7: "总结与下一步",
}

# ── 设计师色板 ─────────────────────────────────────────
INK = RGBColor(0x0A, 0x0A, 0x0B)
PAPER = RGBColor(0xF1, 0xEF, 0xEA)
ACCENT = RGBColor(0x8B, 0x75, 0x6D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1A, 0x1A, 0x1A)


def generate_pptx(
    pages: list[dict[str, Any]],
    project_name: str = "",
    space_type: str = "residential",
    company_name: str = "ARK Design",
    output_path: str | Path = "proposal.pptx",
) -> Path:
    """从 pages 数组生成 PPTX 文件。

    Args:
        pages: ARK narrator 输出的页面数组
        project_name: 项目名称
        space_type: 空间类型
        company_name: 公司名（来自设计师 Profile）
        output_path: 输出路径

    Returns:
        生成的 PPTX 文件路径
    """
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # blank layout

    # ── 封面 ──────────────────────────────────────────
    _add_cover_slide(prs, blank_layout, project_name, space_type, company_name)

    # ── Section → pages grouping ──────────────────────
    section_pages: dict[int, list[dict]] = {}
    for p in pages:
        sec = int(p.get("section", 1))
        section_pages.setdefault(sec, []).append(p)

    # ── 逐 Section 生成 ───────────────────────────────
    for sec in sorted(section_pages):
        sec_name = SECTION_NAMES.get(sec, f"章节 {sec}")
        _add_section_divider(prs, blank_layout, sec, sec_name, company_name)

        for page in section_pages[sec]:
            _add_content_slide(prs, blank_layout, page, sec_name, company_name)

    # ── 结尾页 ────────────────────────────────────────
    _add_closing_slide(prs, blank_layout, company_name)

    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    logger.info("PPTX 已生成: %s (%d 页)", output_path, len(prs.slides))
    return output_path


# ════════════════════════════════════════════════════════
# Slide builders
# ════════════════════════════════════════════════════════


def _add_cover_slide(prs, layout, project_name, space_type, company_name):
    slide = prs.slides.add_slide(layout)
    _set_slide_bg(slide, INK)

    # Chrome bar
    _add_textbox(
        slide, Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4),
        f"{company_name}  ·  提案", size=Pt(10), color=PAPER, bold=False,
    )

    # Title
    space_label = {
        "residential": "私宅", "restaurant": "餐饮", "hotel": "酒店民宿",
        "exhibition": "展厅", "retail": "服务门店",
    }.get(space_type, "设计")

    _add_textbox(
        slide, Inches(0.8), Inches(2.0), Inches(11.7), Inches(1.2),
        f"{space_label} · 概念设计提案", size=Pt(14), color=ACCENT, bold=False,
    )
    _add_textbox(
        slide, Inches(0.8), Inches(3.2), Inches(11.7), Inches(1.6),
        project_name or "设计提案", size=Pt(48), color=WHITE, bold=True,
    )

    # Footer
    _add_textbox(
        slide, Inches(0.8), Inches(6.6), Inches(11.7), Inches(0.4),
        f"{company_name}  ·  {project_name}", size=Pt(9), color=ACCENT,
    )


def _add_section_divider(prs, layout, sec_num, sec_name, company_name):
    slide = prs.slides.add_slide(layout)
    _set_slide_bg(slide, INK)

    _add_textbox(
        slide, Inches(0.8), Inches(2.4), Inches(11.7), Inches(0.6),
        f"0{sec_num}", size=Pt(72), color=ACCENT, bold=True,
    )
    _add_textbox(
        slide, Inches(0.8), Inches(4.0), Inches(11.7), Inches(1.0),
        sec_name, size=Pt(36), color=WHITE, bold=True,
    )
    _add_textbox(
        slide, Inches(0.8), Inches(6.6), Inches(11.7), Inches(0.4),
        f"{company_name}", size=Pt(9), color=ACCENT,
    )


def _add_content_slide(prs, layout, page, sec_name, company_name):
    slide = prs.slides.add_slide(layout)
    _set_slide_bg(slide, PAPER)

    title = page.get("title", "")
    page_type = page.get("pageType", "")
    kicker = page.get("kicker", "")
    blocks = page.get("blocks", [])
    li_note = page.get("liNote", "")

    # Chrome
    _add_textbox(
        slide, Inches(0.6), Inches(0.3), Inches(12.1), Inches(0.35),
        f"{company_name}  ·  {sec_name}", size=Pt(8), color=ACCENT,
    )

    # Kicker + Title
    y = Inches(1.0)
    if kicker:
        _add_textbox(
            slide, Inches(0.8), y, Inches(11.7), Inches(0.4),
            kicker, size=Pt(11), color=ACCENT, bold=False,
        )
        y += Inches(0.45)

    _add_textbox(
        slide, Inches(0.8), y, Inches(11.7), Inches(0.8),
        title, size=Pt(28), color=INK, bold=True,
    )
    y += Inches(0.9)

    # Divider line
    _add_line(slide, Inches(0.8), y, Inches(2.0), y, ACCENT)
    y += Inches(0.3)

    # Blocks
    if blocks:
        # Decide layout based on block count
        if len(blocks) <= 3:
            y = _render_pillar_blocks(slide, blocks, y)
        elif len(blocks) <= 6:
            y = _render_step_blocks(slide, blocks, y)
        else:
            y = _render_grid_blocks(slide, blocks, y)

    # liNote (designer comment) at bottom
    if li_note:
        _add_textbox(
            slide, Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.5),
            f"💬 {li_note}", size=Pt(10), color=ACCENT,
        )


def _add_closing_slide(prs, layout, company_name):
    slide = prs.slides.add_slide(layout)
    _set_slide_bg(slide, INK)

    _add_textbox(
        slide, Inches(0.8), Inches(2.8), Inches(11.7), Inches(1.0),
        "感谢阅读", size=Pt(48), color=WHITE, bold=True,
    )
    _add_textbox(
        slide, Inches(0.8), Inches(4.0), Inches(11.7), Inches(0.6),
        "期待与您进一步探讨方案细节", size=Pt(18), color=ACCENT,
    )
    _add_textbox(
        slide, Inches(0.8), Inches(6.6), Inches(11.7), Inches(0.4),
        f"本提案由 {company_name} AI 设计引擎生成", size=Pt(9), color=ACCENT,
    )


# ════════════════════════════════════════════════════════
# Block renderers
# ════════════════════════════════════════════════════════

def _render_pillar_blocks(slide, blocks, y):
    """3 列支柱布局"""
    w = Inches(3.7)
    gap = Inches(0.3)
    x_start = Inches(0.8)

    for i, block in enumerate(blocks[:3]):
        x = x_start + i * (w + gap)
        _add_textbox(
            slide, x, y, w, Inches(0.4),
            block.get("title", "") or block.get("label", ""),
            size=Pt(14), color=INK, bold=True,
        )
        desc = block.get("desc", "") or block.get("value", "") or ""
        _add_textbox(
            slide, x, y + Inches(0.4), w, Inches(1.8),
            str(desc), size=Pt(11), color=RGBColor(0x55, 0x55, 0x55),
        )

    return y + Inches(2.6)


def _render_step_blocks(slide, blocks, y):
    """步骤/流水线布局"""
    for i, block in enumerate(blocks[:6]):
        _add_textbox(
            slide, Inches(1.2), y, Inches(0.5), Inches(0.4),
            f"{i + 1:02d}", size=Pt(18), color=ACCENT, bold=True,
        )
        _add_textbox(
            slide, Inches(1.9), y, Inches(9.0), Inches(0.35),
            block.get("title", "") or block.get("label", ""),
            size=Pt(14), color=INK, bold=True,
        )
        desc = block.get("desc", "") or block.get("value", "") or ""
        if desc:
            _add_textbox(
                slide, Inches(1.9), y + Inches(0.35), Inches(9.0), Inches(0.6),
                str(desc), size=Pt(11), color=RGBColor(0x55, 0x55, 0x55),
            )
        y += Inches(1.0)
    return y


def _render_grid_blocks(slide, blocks, y):
    """网格布局 (≥7 blocks)"""
    cols = 3
    w = Inches(3.7)
    h = Inches(2.0)
    gap_x = Inches(0.3)
    gap_y = Inches(0.15)
    x_start = Inches(0.8)

    for i, block in enumerate(blocks):
        col = i % cols
        row = i // cols
        x = x_start + col * (w + gap_x)
        yy = y + row * (h + gap_y)

        _add_textbox(
            slide, x, yy, w, Inches(0.4),
            block.get("title", "") or block.get("label", ""),
            size=Pt(13), color=INK, bold=True,
        )
        desc = block.get("desc", "") or block.get("value", "") or ""
        _add_textbox(
            slide, x, yy + Inches(0.4), w, Inches(1.4),
            str(desc), size=Pt(10), color=RGBColor(0x55, 0x55, 0x55),
        )

    return y + Inches(4.5)


# ════════════════════════════════════════════════════════
# Drawing helpers
# ════════════════════════════════════════════════════════

def _set_slide_bg(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_textbox(slide, left, top, width, height, text, *, size=Pt(12), color=INK, bold=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = str(text)
    p.font.size = size
    p.font.color.rgb = color
    p.font.bold = bold
    return txBox


def _add_line(slide, x1, y1, x2, y2, color, width=Pt(2)):
    connector = slide.shapes.add_connector(
        1, x1, y1, x2, y2  # MSO_CONNECTOR.STRAIGHT = 1
    )
    connector.line.color.rgb = color
    connector.line.width = width
